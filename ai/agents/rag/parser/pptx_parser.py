import io
from pptx import Presentation
from .parser_interface import BaseParser

class PptxParser(BaseParser):
    def parse(self, file_bytes: bytes, file_name: str) -> str:
        try:
            prs = Presentation(io.BytesIO(file_bytes))
            output = []
            for i, slide in enumerate(prs.slides):
                output.append(f"## Slide {i+1}")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        output.append(shape.text.strip())
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        output.append(f"Notes: {notes}")
                output.append("")
            return "\n".join(output)
        except Exception as e:
            return f"Failed to parse PowerPoint file {file_name}: {str(e)}"
